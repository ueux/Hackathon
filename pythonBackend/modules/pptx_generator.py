# modules/pptx_generator.py - Fixed version with proper placeholder handling
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import os

class PPTXGenerator:
    def __init__(self, output_file, template=None):
        self.output_file = output_file
        self.presentation = self._init_presentation(template)
        self._setup_styles()
        self.slide_count = 0
        self.max_slides = 12

    def _init_presentation(self, template):
        """Initialize presentation with proper error handling"""
        try:
            if template and os.path.exists(template):
                return Presentation(template)
            return Presentation()
        except Exception as e:
            raise Exception(f"Failed to initialize presentation: {str(e)}")

    def _setup_styles(self):
        """Define presentation styles"""
        self.styles = {
            'title': {
                'font_size': Pt(36),
                'font_name': 'Calibri Light',
                'color': RGBColor(23, 54, 93),
                'alignment': PP_ALIGN.LEFT
            },
            'content': {
                'font_size': Pt(20),
                'font_name': 'Calibri',
                'color': RGBColor(50, 50, 50),
                'line_spacing': 1.3
            }
        }

    def _safe_get_placeholder(self, slide, idx):
        """Safely get placeholder by index"""
        try:
            if len(slide.placeholders) > idx:
                return slide.placeholders[idx]
            return None
        except Exception:
            return None

    def generate_presentation(self, structured_content):
        """Generate the optimized presentation"""
        if not structured_content:
            raise ValueError("No content provided")

        # Add title slide
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title = self._safe_get_placeholder(title_slide, 0)
        subtitle = self._safe_get_placeholder(title_slide, 1)

        if title:
            title.text = "Executive Summary"
            self._format_text(title, 'title')
        if subtitle:
            subtitle.text = "Condensed Key Points"
            self._format_text(subtitle, 'content')

        self.slide_count += 1

        # Process content
        for section in structured_content[:8]:  # Limit to 8 sections
            if self.slide_count >= self.max_slides:
                break
            self._add_section(section)

        # Add closing slide
        self._add_closing_slide()
        self.presentation.save(self.output_file)

    def _add_section(self, section):
        """Add a content section"""
        # Section title slide
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._safe_get_placeholder(title_slide, 0)
        if title:
            title.text = section["title"]
            self._format_text(title, 'title')
        self.slide_count += 1

        # Content slides
        content = self._condense_content(section["content"])
        for i in range(0, len(content), 4):  # 4 bullets per slide
            if self.slide_count >= self.max_slides:
                return
            self._add_content_slide(content[i:i+4])
            self.slide_count += 1

    def _add_content_slide(self, bullets):
        """Add a content slide with bullet points"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        content = self._safe_get_placeholder(slide, 1)

        if content:
            tf = content.text_frame
            tf.clear()
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                self._format_text(p, 'content')

    def _condense_content(self, content):
        """Condense content to key points"""
        return [item for item in content if item.strip()][:16]  # Max 16 bullets total

    def _add_closing_slide(self):
        """Add final summary slide"""
        if self.slide_count >= self.max_slides:
            return

        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._safe_get_placeholder(slide, 0)
        content = self._safe_get_placeholder(slide, 1)

        if title:
            title.text = "Key Takeaways"
            self._format_text(title, 'title')
        if content:
            tf = content.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = "Thank you for your attention"
            self._format_text(p, 'content')

    def _format_text(self, element, style_type):
        """Apply formatting to text elements"""
        style = self.styles[style_type]
        if hasattr(element, 'text_frame'):
            for paragraph in element.text_frame.paragraphs:
                paragraph.alignment = style.get('alignment', PP_ALIGN.LEFT)
                for run in paragraph.runs:
                    run.font.size = style['font_size']
                    run.font.name = style['font_name']
                    run.font.color.rgb = style['color']
        elif hasattr(element, 'runs'):
            for run in element.runs:
                run.font.size = style['font_size']
                run.font.name = style['font_name']
                run.font.color.rgb = style['color']