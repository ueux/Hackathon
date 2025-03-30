# modules/data_analyzer.py
import re
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from collections import defaultdict
from pptx.util import Inches
import tempfile
import os

class DataAnalyzer:
    # [Previous methods remain the same until add_to_presentation]
    def __init__(self):
        self.numeric_data = defaultdict(list)
        self.tables = []
        self.figures = []

    def extract_and_analyze(self, content):
        """Main method to extract and analyze numerical data"""
        self._extract_numerical_data(content)

    def _extract_numerical_data(self, content):
        """Extract numerical data from document content"""
        num_pattern = r"""
            (?:^|\s)                         # Start of string or whitespace
            (?:\$?\s*-?\d{1,3}(?:,\d{3})*   # Numbers with commas
            (?:\.\d+)?\s*%?                  # Decimals and percentages
            |\.\d+\s*%?)                     # Decimal numbers
            (?:\s*[-–]\s*                    # Range separator
            \$?\s*-?\d{1,3}(?:,\d{3})*      # Second number in range
            (?:\.\d+)?\s*%?)?                # Optional decimal/percentage
            (?=\s|$|[^\w.-])                 # Lookahead for boundary
        """

        for section in content:
            section_data = {
                'title': section['title'],
                'values': [],
                'units': [],
                'context': []
            }

            for line in section['content']:
                matches = re.finditer(num_pattern, line, re.VERBOSE)
                for match in matches:
                    value_str = match.group().strip()

            if section_data['values']:
                self.numeric_data[section['title']] = section_data

    def add_to_presentation(self, presentation):
        """Add analysis results to PowerPoint presentation"""
        # Create temp directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save all figures as PNG images
            img_paths = []
            for i, fig in enumerate(self.figures):
                img_path = os.path.join(temp_dir, f"plot_{i}.png")
                fig.savefig(img_path, dpi=300, bbox_inches='tight')
                img_paths.append(img_path)
                plt.close(fig)

            # Add tables first
            self._add_tables_to_presentation(presentation)

            # Add images to slides
            for i, img_path in enumerate(img_paths):
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Title only
                slide.shapes.title.text = f"Data Visualization {i+1}"

                # Add image to slide (centered)
                left = (presentation.slide_width - Inches(7)) // 2
                slide.shapes.add_picture(
                    img_path,
                    left, Inches(1.5),
                    width=Inches(7),
                    height=Inches(5)
                )

    def _add_tables_to_presentation(self, presentation):
        """Helper method to add tables to slides"""
        for table in self.tables:
            # Add data table slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = table['title']

            # Create table
            rows = len(table['data']) + 1  # +1 for header
            cols = len(table['data'].columns)
            left, top = Inches(0.5), Inches(1.5)
            width, height = Inches(9), Inches(0.3 * rows)

            pptx_table = slide.shapes.add_table(
                rows, cols, left, top, width, height
            ).table

            # Add headers
            for col_idx, col_name in enumerate(table['data'].columns):
                cell = pptx_table.cell(0, col_idx)
                cell.text = str(col_name)
                cell.text_frame.paragraphs[0].font.bold = True

            # Add data rows
            for row_idx, (_, row) in enumerate(table['data'].iterrows(), 1):
                for col_idx, col_name in enumerate(table['data'].columns):
                    pptx_table.cell(row_idx, col_idx).text = str(row[col_name])

            # Add statistics slide
            self._add_stats_slide(presentation, table)

    def _add_stats_slide(self, presentation, table):
        """Add statistics summary slide"""
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Statistics: {table['title']}"

        # Create stats table
        stats = table['stats']
        rows = len(stats) + 1
        left, top = Inches(2), Inches(1.5)
        width, height = Inches(4), Inches(0.3 * rows)

        pptx_table = slide.shapes.add_table(
            rows, 2, left, top, width, height
        ).table

        # Add headers
        pptx_table.cell(0, 0).text = "Metric"
        pptx_table.cell(0, 1).text = "Value"

        # Add stats rows
        for row_idx, (stat, val) in enumerate(stats.items(), 1):
            pptx_table.cell(row_idx, 0).text = stat.replace('_', ' ').title()
            pptx_table.cell(row_idx, 1).text = f"{val:.2f}" if isinstance(val, (float, np.floating)) else str(val)