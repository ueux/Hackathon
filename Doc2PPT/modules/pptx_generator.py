from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List, Dict
import os

class PPTXGenerator:
    def __init__(self, output_path: str, template_path: str = None):
        self.output_path = output_path
        self.presentation = self._init_presentation(template_path)
        self._setup_styles()
        self.slide_count = 0
        self.max_slides = 20  # Default maximum slides

    def generate_presentation(self, content: List[Dict],
                           include_summary: bool = True,
                           include_appendix: bool = False,
                           audience_level: str = 'executive',
                           presentation_length: str = 'medium'):
        """
        Generate presentation with content and customization options

        Args:
            content: Structured content from document parser
            include_summary: Whether to include summary slide
            include_appendix: Whether to include appendix
            audience_level: 'executive', 'management', or 'technical'
            presentation_length: 'short', 'medium', or 'long'
        """
        try:
            # Set slide limits based on length preference
            self._set_slide_limits(presentation_length)

            # Add title slide
            self._add_title_slide("Document Conversion", "Automatically generated presentation")

            # Add summary if requested
            if include_summary:
                self._add_summary_slide(content)

            # Add content slides filtered by audience level
            filtered_content = self._filter_by_audience(content, audience_level)
            for section in filtered_content:
                if self.slide_count >= self.max_slides:
                    break
                self._add_section(section)

            # Add appendix if requested
            if include_appendix and self.slide_count < self.max_slides:
                self._add_appendix_slide()

            # Add closing slide
            if self.slide_count < self.max_slides:
                self._add_closing_slide()

            # Save the final presentation
            self.presentation.save(self.output_path)
            return True

        except Exception as e:
            raise Exception(f"Failed to generate presentation: {str(e)}")

    def _init_presentation(self, template_path: str = None):
        """Initialize presentation with proper error handling"""
        try:
            if template_path and os.path.exists(template_path):
                return Presentation(template_path)
            return Presentation()
        except Exception as e:
            raise Exception(f"Failed to initialize presentation: {str(e)}")

    def _setup_styles(self):
        """Define presentation styles"""
        self.styles = {
            'title': {
                'font_size': Pt(36),
                'font_name': 'Calibri Light',
                'color': RGBColor(23, 54, 93),
                'alignment': PP_ALIGN.LEFT
            },
            'content': {
                'font_size': Pt(20),
                'font_name': 'Calibri',
                'color': RGBColor(50, 50, 50),
                'line_spacing': 1.3
            },
            'executive': {
                'font_size': Pt(28),
                'max_bullets': 3
            },
            'management': {
                'font_size': Pt(24),
                'max_bullets': 5
            },
            'technical': {
                'font_size': Pt(20),
                'max_bullets': 7
            }
        }

    def _set_slide_limits(self, length: str):
        """Set maximum slides based on length preference"""
        self.max_slides = {
            'short': 10,
            'medium': 15,
            'long': 20
        }.get(length.lower(), 15)

    def _filter_by_audience(self, content: List[Dict], audience_level: str) -> List[Dict]:
        """Filter content based on audience level"""
        if audience_level == 'executive':
            return [section for section in content
                   if any(kw in section['title'].lower()
                   for kw in ['summary', 'key', 'result', 'conclusion'])]
        elif audience_level == 'management':
            return content[:10]  # First 10 sections for management
        return content  # All content for technical audience

    def _add_title_slide(self, title: str, subtitle: str):
        """Add title slide to presentation"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle
        self._format_text(title_shape, 'title')
        self._format_text(subtitle_shape, 'content')
        self.slide_count += 1

    def _add_summary_slide(self, content: List[Dict]):
        """Add summary slide with key points"""
        if self.slide_count >= self.max_slides:
            return

        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._get_placeholder(slide, 0)
        content_ph = self._get_placeholder(slide, 1)

        if title:
            title.text = "Document Summary"
            self._format_text(title, 'title')

        if content_ph:
            tf = content_ph.text_frame
            tf.clear()

            # Add key points from first few sections
            for section in content[:3]:
                p = tf.add_paragraph()
                p.text = section['title']
                p.level = 0
                self._format_text(p, 'content')

                for point in section['content'][:2]:
                    p = tf.add_paragraph()
                    p.text = f"- {point}"
                    p.level = 1
                    self._format_text(p, 'content')

        self.slide_count += 1

    def _add_section(self, section: Dict):
        """Add a content section with title and bullet points"""
        # Section title slide
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._get_placeholder(title_slide, 0)
        if title:
            title.text = section["title"]
            self._format_text(title, 'title')
        self.slide_count += 1

        # Content slides with bullet points
        max_bullets = self.styles.get('content', {}).get('max_bullets', 5)
        for i in range(0, len(section['content']), max_bullets):
            if self.slide_count >= self.max_slides:
                return
            self._add_bullet_slide(
                section['title'],
                section['content'][i:i+max_bullets]
            )
            self.slide_count += 1

    def _add_bullet_slide(self, section_title: str, bullets: List[str]):
        """Add slide with bullet points"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._get_placeholder(slide, 0)
        content = self._get_placeholder(slide, 1)

        if title:
            title.text = section_title
            self._format_text(title, 'title')

        if content:
            tf = content.text_frame
            tf.clear()
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                self._format_text(p, 'content')

    def _add_appendix_slide(self):
        """Add appendix slide if requested"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._get_placeholder(slide, 0)
        if title:
            title.text = "Appendix"
            self._format_text(title, 'title')

    def _add_closing_slide(self):
        """Add final closing slide"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = self._get_placeholder(slide, 0)
        content = self._get_placeholder(slide, 1)

        if title:
            title.text = "Thank You"
            self._format_text(title, 'title')
        if content:
            content.text = "This presentation was automatically generated"
            self._format_text(content, 'content')

    def _get_placeholder(self, slide, idx: int):
        """Safely get placeholder by index"""
        try:
            if len(slide.placeholders) > idx:
                return slide.placeholders[idx]
            return None
        except Exception:
            return None

    def _format_text(self, element, style_type: str):
        """Apply consistent formatting to text elements"""
        style = self.styles.get(style_type, {})
        if not style:
            return

        if hasattr(element, 'text_frame'):
            for paragraph in element.text_frame.paragraphs:
                if 'alignment' in style:
                    paragraph.alignment = style['alignment']
                for run in paragraph.runs:
                    self._apply_font_style(run, style)
        elif hasattr(element, 'runs'):
            for run in element.runs:
                self._apply_font_style(run, style)
        elif hasattr(element, 'font'):
            self._apply_font_style(element, style)

    def _apply_font_style(self, run, style: Dict):
        """Apply font styles to a run"""
        if 'font_size' in style:
            run.font.size = style['font_size']
        if 'font_name' in style:
            run.font.name = style['font_name']
        if 'color' in style:
            run.font.color.rgb = style['color']